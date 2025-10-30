# translation_engine.py
# Core translation logic extracted from PPT_Language_Change.py

import os
import re
import time
from pptx import Presentation
from pptx.dml.color import RGBColor
import openai


# ---------- [서식 보존을 위한 태깅/복원 유틸] ----------
RUN_TAG = re.compile(r"\[\[R(\d+)\]\]|\[\[/R(\d+)\]\]")

# ==== [설정] =====================================================
SORRY_PATTERNS = [
    "i'm sorry", "i am sorry", "sorry, but", "apologize",
    "죄송하지만", "죄송합니다", "번역할 내용이 없습니다"
]

LANG_OPTIONS = [
    "English",
    "Indonesian",
    "Italian",
    "French",
    "Spanish",
    "Korean",
    "Japanese",
    "Russian",
    "German",
    "Portuguese",
    "Chinese (Simplified)",
    "Chinese (Traditional)",
]

TONE_OPTIONS = [
    "기본값",
    "Med/Pharma Pro (20y)",   # 의료기기/전문약사 20년 전문가
    "Medical/Science Expert", # 의학/과학 번역 전문가
    "Beauty Pro (20y, chic)", # 세련된 뷰티 20년 전문가
    "GenZ Female (20s)",      # 20대 여성 타깃
]

OPENAI_MODEL = "gpt-5"
DEEPSEEK_MODEL = "deepseek-chat"
SLEEP_SEC = 0


def is_effectively_empty_tagged(tagged_text: str) -> bool:
    """[[R#]]마커를 제거하고 남는 콘텐츠가 실질적으로 비었는지 판단"""
    stripped = RUN_TAG.sub("", tagged_text)  # 마커 제거
    return stripped.strip() == ""  # 공백만 남으면 빈 것으로 간주


def looks_like_apology(text: str) -> bool:
    low = (text or "").lower()
    return any(p in low for p in SORRY_PATTERNS)


def unique_path(path: str) -> str:  
    if not os.path.exists(path):
        return path
    base, ext = os.path.splitext(path)
    i = 1
    while True:
        candidate = f"{base} ({i}){ext}"
        if not os.path.exists(candidate):
            return candidate
        i += 1


def create_openai_client(api_key: str):
    """OpenAI 클라이언트 생성"""
    return openai.OpenAI(api_key=api_key)


def create_deepseek_client(api_key: str):
    """DeepSeek 클라이언트 생성"""
    return openai.OpenAI(
        api_key=api_key,
        base_url="https://api.deepseek.com"
    )


def safe_request(client, prompt, retries=3, delay=3, use_deepseek=False):
    for attempt in range(retries):
        try:
            model = DEEPSEEK_MODEL if use_deepseek else OPENAI_MODEL
            resp = client.chat.completions.create(
                model=model,
                messages=[{"role": "user", "content": prompt}],
                timeout=60,
            )
            content = ""
            if resp and hasattr(resp, "choices") and resp.choices:
                content = getattr(resp.choices[0].message, "content", "") or ""
            if content:
                return content.strip()
        except Exception as e:
            print(f"⚠️ Error (attempt {attempt+1}): {e}")
            with open("error.log", "a", encoding="utf-8") as f:
                f.write(f"[Attempt {attempt+1}] {e}\n")
            time.sleep(delay)
    return ""


def _font_attrs(run):
    f = run.font
    # 값 그대로 보존(None 포함)
    name = f.name                 # None이면 테마/마스터 상속
    size = f.size.pt if f.size else None
    bold = f.bold                 # True/False/None
    italic = f.italic             # True/False/None
    underline = f.underline       # True/False/None

    rgb = None
    if f.color is not None and getattr(f.color, "rgb", None) is not None:
        rgb = (f.color.rgb[0], f.color.rgb[1], f.color.rgb[2])

    return {"name": name, "size": size, "bold": bold, "italic": italic,
            "underline": underline, "rgb": rgb}


def _apply_font_attrs(run, attrs):
    from pptx.util import Pt
    f = run.font

    if attrs.get("name") is not None:
        f.name = attrs["name"]
    if attrs.get("size") is not None:
        f.size = Pt(attrs["size"])
    if attrs.get("bold") is not None:
        f.bold = attrs["bold"]
    if attrs.get("italic") is not None:
        f.italic = attrs["italic"]
    if attrs.get("underline") is not None:
        f.underline = attrs["underline"]
    if attrs.get("rgb") is not None:
        r, g, b = attrs["rgb"]
        f.color.rgb = RGBColor(r, g, b)


def adjust_font_size_for_translation(original_size, font_scale):
    """번역된 텍스트의 폰트 크기를 조정"""
    if original_size is None:
        return None
    return original_size * font_scale


def _apply_font_attrs_with_scale(run, attrs, font_scale=1.0):
    """폰트 속성을 적용하되 크기를 조정"""
    from pptx.util import Pt
    f = run.font

    if attrs.get("name") is not None:
        f.name = attrs["name"]
    
    # 폰트 크기 조정 적용
    if attrs.get("size") is not None:
        adjusted_size = adjust_font_size_for_translation(attrs["size"], font_scale)
        f.size = Pt(adjusted_size)
    
    if attrs.get("bold") is not None:
        f.bold = attrs["bold"]
    if attrs.get("italic") is not None:
        f.italic = attrs["italic"]
    if attrs.get("underline") is not None:
        f.underline = attrs["underline"]
    if attrs.get("rgb") is not None:
        r, g, b = attrs["rgb"]
        f.color.rgb = RGBColor(r, g, b)


def tag_paragraph(paragraph):
    text_parts, style_map, idx = [], {}, 1
    for run in paragraph.runs:
        t = run.text or ""
        if not t:
            continue
        style_map[idx] = _font_attrs(run)
        text_parts.append(f"[[R{idx}]]{t}[[/R{idx}]]")
        idx += 1
    return "".join(text_parts), style_map


def rebuild_paragraph_from_tagged(paragraph, translated, style_map, font_scale=1.0):
    while paragraph.runs:
        paragraph.runs[0]._r.getparent().remove(paragraph.runs[0]._r)

    tokens, pos = [], 0
    for m in RUN_TAG.finditer(translated):
        s, e = m.span()
        if s > pos:
            tokens.append(("text", translated[pos:s]))
        if m.group(1):
            tokens.append(("start", int(m.group(1))))
        if m.group(2):
            tokens.append(("end", int(m.group(2))))
        pos = e
    if pos < len(translated):
        tokens.append(("text", translated[pos:]))

    out_runs, stack, buffer = [], [], {}
    for ttype, value in tokens:
        if ttype == "start":
            stack.append(value)
            buffer.setdefault(value, [])
        elif ttype == "end":
            if stack and stack[-1] == value:
                stack.pop()
            joined = "".join(buffer.get(value, []))
            out_runs.append((value, joined))
            buffer[value] = []
        else:
            if stack:
                buffer[stack[-1]].append(value)
            else:
                out_runs.append((None, value))

    for run_id, buf in buffer.items():
        if buf:
            out_runs.append((run_id, "".join(buf)))

    if not out_runs:
        r = paragraph.add_run()
        r.text = translated
        return

    for run_id, txt in out_runs:
        if not txt:
            continue
        r = paragraph.add_run()
        r.text = txt
        if run_id and run_id in style_map:
            if font_scale != 1.0:
                _apply_font_attrs_with_scale(r, style_map[run_id], font_scale)
            else:
                _apply_font_attrs(r, style_map[run_id])


def _parse_run_chunks(translated):
    # R# → 텍스트 매핑과, 마커 밖 텍스트 존재 여부를 판정
    ids = []
    chunks = {}
    stack = []
    buf = {}
    outside = []

    pos = 0
    for m in RUN_TAG.finditer(translated):
        s, e = m.span()
        if s > pos:
            if stack:
                buf.setdefault(stack[-1], []).append(translated[pos:s])
            else:
                outside.append(translated[pos:s])
        if m.group(1):  # [[R#]]
            rid = int(m.group(1)); stack.append(rid); ids.append(rid)
        if m.group(2):  # [[/R#]]
            rid = int(m.group(2))
            if stack and stack[-1] == rid:
                stack.pop()
                joined = "".join(buf.get(rid, []))
                chunks[rid] = joined
                buf[rid] = []
        pos = e
    if pos < len(translated):
        if stack:
            buf.setdefault(stack[-1], []).append(translated[pos:])
        else:
            outside.append(translated[pos:])

    # 닫히지 않은 버퍼 처리
    for rid, lst in buf.items():
        if lst:
            chunks[rid] = chunks.get(rid, "") + "".join(lst)

    has_outside = any(t.strip() for t in outside)
    return ids, chunks, has_outside


def try_inplace_update_paragraph(paragraph, translated, font_scale=1.0):
    """마커가 1..N으로 정확히 존재하고, 마커 밖 텍스트가 없으면
    기존 runs에 텍스트만 주입하여 서식을 100% 유지한다."""
    ids, chunks, has_outside = _parse_run_chunks(translated)
    runs = [r for r in paragraph.runs if (r.text or "") != ""]
    N = len(runs)

    # 조건: 마커 밖 텍스트가 없어야 하고, R1..RN이 정확히 한 번씩 존재
    if has_outside or N == 0 or set(ids) != set(range(1, N+1)) or any(ids.count(i) != 1 for i in range(1, N+1)):
        return False

    for i, run in enumerate(runs, start=1):
        run.text = chunks.get(i, "")
        # 폰트 크기 조정 적용
        if font_scale != 1.0 and run.font.size is not None:
            from pptx.util import Pt
            original_size = run.font.size.pt
            adjusted_size = original_size * font_scale
            run.font.size = Pt(adjusted_size)
    return True


# ---------- [프롬프트 빌더] ----------
def build_tone_instructions(tone: str) -> str:
    """
    선택한 톤에 맞는 지시문을 반환
    """
    if tone == "기본값":
        return (
            "Use a natural, professional beauty-industry tone localized to the target market. "
            "Keep terminology consistent with beauty marketing and professional skincare. "
            "Be clear and persuasive without hype; avoid overpromising."
        )
    if tone == "Med/Pharma Pro (20y)":
        return (
            "Use a formal, clinically precise B2B tone suitable for medical devices and professional pharmacists. "
            "Prioritize clarity, compliance, and evidence-based wording. Avoid hype. "
            "Prefer terminology used in regulatory, clinical, and professional settings."
        )
    if tone == "Beauty Pro (20y, chic)":
        return (
            "Use a refined, polished professional tone common in premium beauty and aesthetic clinics. "
            "Balance expertise with approachable elegance. Maintain brand voice consistency without overpromising."
        )
    if tone == "GenZ Female (20s)":
        return (
            "Use a modern, friendly, and concise tone tailored for women in their 20s. "
            "Be clear and engaging for social content, but avoid slang overload, emojis, and exaggerated claims."
        )
    if tone == "Medical/Science Expert":
        return (
            "Use a precise, scientific, and clinically accurate tone suitable for medical and scientific documentation. "
            "Maintain strict adherence to medical terminology, biological processes, and scientific accuracy. "
            "Preserve technical terms, Latin names, and scientific nomenclature exactly as they appear. "
            "Ensure translations are suitable for medical professionals, researchers, and regulatory submissions."
        )
    # fallback
    return "Use a neutral professional tone appropriate for the beauty industry."


def build_chinese_prompt(tagged_text: str, target_lang: str) -> str:
    """
    전문적인 한국어→중국어 번역을 위한 프롬프트 (간체/번체 구분)
    """
    chinese_type = "간체" if "Simplified" in target_lang else "번체"
    
    return (
        f"당신은 한국어를 정확하고 자연스러운 중국어({chinese_type})로 번역하는 전문가입니다.\n\n"
        f"# 주요 특징\n"
        f"- 프레젠테이션, 보고서, 비즈니스 문서 등에 적합한 공식적이고 세련된 표현 사용\n"
        f"- 문맥을 고려하여 문장의 의미와 뉘앙스를 세밀하게 분석하여 적절한 표현으로 번역\n"
        f"- 원문의 의도와 어조를 유지하되, 중국 원어민이 자연스럽게 들리도록 문장 구조 조정\n"
        f"- 번역 이외의 불필요한 설명 금지\n"
        f"- 창의적 재해석 없이 원문에 충실한 번역 수행\n"
        f"- 반드시 중국어 {chinese_type}로 번역하세요\n\n"
        f"# 고유명사 처리 규칙\n"
        f"- '피더린'은 'PYDERIN'으로 번역하세요 (브랜드명이므로 대문자로)\n"
        f"- 기타 고유명사(인명, 지명, 회사명, 브랜드명 등)는 번역하지 말고 원문 그대로 유지하세요\n"
        f"- 영어 고유명사는 그대로 유지하세요\n\n"
        f"다음 한국어 텍스트를 자연스러운 중국어({chinese_type})로 번역하세요.\n"
        f"중요: [[R1]]...[[/R1]] 같은 마커 태그는 절대 변경하거나 제거하지 마세요. 정확히 그대로 유지해야 합니다.\n\n"
        f"번역할 텍스트:\n{tagged_text}"
    )


def build_medical_science_prompt(tagged_text: str, target_lang: str) -> str:
    """
    의학/과학 전문 번역을 위한 프롬프트
    """
    return (
        f"You are a medical and scientific translation expert with 20+ years of experience. "
        f"Translate the following Korean medical/scientific text into precise, professional {target_lang}. "
        f"Only return the translated text. If there is nothing to translate, return an empty string.\n\n"
        f"# Translation Guidelines:\n"
        f"- Maintain strict scientific accuracy and medical terminology\n"
        f"- Preserve all technical terms, Latin names, and scientific nomenclature exactly\n"
        f"- Use appropriate medical/scientific vocabulary for the target language\n"
        f"- Ensure the translation is suitable for medical professionals and researchers\n"
        f"- Maintain the formal, clinical tone of scientific documentation\n"
        f"- Keep biological processes, anatomical terms, and chemical names precise\n"
        f"- If the source is already in {target_lang}, lightly copyedit for scientific accuracy\n\n"
        f"# Critical Requirements:\n"
        f"- Do NOT alter or remove any marker tags like [[R1]]...[[/R1]]\n"
        f"- Keep all markers exactly as-is and in correct pairs\n"
        f"- Preserve the exact structure and formatting\n\n"
        f"Translate the following text:\n{tagged_text}"
    )


def build_prompt(tagged_text: str, target_lang: str, tone: str) -> str:
    # Medical/Science Expert uses specialized prompt
    if tone == "Medical/Science Expert":
        return build_medical_science_prompt(tagged_text, target_lang)
    
    # Chinese translation uses specialized prompt
    if "Chinese" in target_lang:
        return build_chinese_prompt(tagged_text, target_lang)
    
    tone_text = build_tone_instructions(tone)
    return (
        f"Translate the following beauty industry presentation text into natural, professional {target_lang}. "
        f"Only return the translated text. If there is nothing to translate, return an empty string. "
        f"Context: {tone_text} "
        f"Avoid literal translation—use expressions that sound natural for beauty marketing and professional skincare. "
        f"If the source is already in {target_lang}, lightly copyedit for clarity, consistency, and terminology. "
        f"CRITICAL: Do NOT alter or remove any marker tags like [[R1]]...[[/R1]]. Keep them exactly as-is and in correct pairs. "
        f"Return the translated text with all markers preserved:\n\n{tagged_text}"
    )


# ---------- [번역 호출] ----------
def gpt_translate_tagged(tagged_text: str, client, target_lang: str, tone: str, use_deepseek=False) -> str:
    # 진짜 내용이 없으면 번역 스킵
    if not tagged_text.strip() or is_effectively_empty_tagged(tagged_text):
        return ""

    # 중국어 번역의 경우 DeepSeek 사용
    if "Chinese" in target_lang and use_deepseek:
        deepseek_client = create_deepseek_client(client.api_key)
        prompt = build_chinese_prompt(tagged_text, target_lang)
        content = safe_request(deepseek_client, prompt, retries=3, delay=3, use_deepseek=True)
    else:
        prompt = build_prompt(tagged_text, target_lang, tone)
        content = safe_request(client, prompt, retries=3, delay=3)

    # 실패 시 원문(마커 포함) 반환 → 원문 유지
    if not content:
        return tagged_text

    # 사과문/에러문구가 들어오면 원문 유지
    if looks_like_apology(content):
        return tagged_text

    return content


def process_nested_shapes(shapes, target_lang, tone, client, use_deepseek, font_scale, progress_callback, slide_num, total_slides, shape_path="", should_stop=None):
    """중첩된 shape들을 재귀적으로 처리"""
    for shape_idx, shape in enumerate(shapes):
        # 중지 신호 확인
        if should_stop and should_stop():
            print("   ⏹️ 번역이 중지되었습니다.")
            return False
            
        current_path = f"{shape_path}.{shape_idx}" if shape_path else str(shape_idx)
        
        # 텍스트 프레임 처리
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            tf = shape.text_frame
            for p_idx, p in enumerate(tf.paragraphs):
                tagged, style_map = tag_paragraph(p)
                if not tagged:
                    continue
                preview = (tagged[:40] + "...") if len(tagged) > 40 else tagged
                print(f"   🔤 번역 중(중첩텍스트): {preview}")
                if progress_callback:
                    progress_callback(slide_num, total_slides, f"중첩 텍스트 번역 중: {preview}")
                
                translated = gpt_translate_tagged(tagged, client, target_lang, tone, use_deepseek)
                translated = translated.strip().strip('"').strip("'")
                
                if not try_inplace_update_paragraph(p, translated, font_scale):
                    rebuild_paragraph_from_tagged(p, translated, style_map, font_scale)
                
                time.sleep(SLEEP_SEC)
        
        # 표 처리 (중첩된 표 포함)
        elif getattr(shape, "has_table", False) and shape.has_table:
            print(f"   📊 표 처리 중 (경로: {current_path})")
            for row_idx, row in enumerate(shape.table.rows):
                for cell_idx, cell in enumerate(row.cells):
                    # 셀의 텍스트 프레임 처리
                    if getattr(cell, "text_frame", None):
                        tf = cell.text_frame
                        for p_idx, p in enumerate(tf.paragraphs):
                            tagged, style_map = tag_paragraph(p)
                            if not tagged:
                                continue
                            translated = gpt_translate_tagged(tagged, client, target_lang, tone, use_deepseek)
                            translated = translated.strip().strip('"').strip("'")
                            
                            if not try_inplace_update_paragraph(p, translated, font_scale):
                                rebuild_paragraph_from_tagged(p, translated, style_map, font_scale)
                            time.sleep(SLEEP_SEC)
                    
                    # 셀 안의 다른 shape들 처리 (중첩된 표, 텍스트박스 등)
                    if hasattr(cell, 'shapes') and cell.shapes:
                        print(f"     🔍 셀 내부 shape 발견 (행:{row_idx}, 열:{cell_idx})")
                        if not process_nested_shapes(
                            cell.shapes, target_lang, tone, client, use_deepseek, 
                            font_scale, progress_callback, slide_num, total_slides, f"{current_path}.table_{row_idx}_{cell_idx}", should_stop
                        ):
                            return False
        
        # 기타 shape 타입들도 확인 (그룹, 텍스트박스 등)
        elif hasattr(shape, 'shapes') and shape.shapes:
            print(f"   🔍 그룹 shape 발견 (경로: {current_path})")
            if not process_nested_shapes(
                shape.shapes, target_lang, tone, client, use_deepseek, 
                font_scale, progress_callback, slide_num, total_slides, current_path, should_stop
            ):
                return False
    
    return True


def translate_presentation(pptx_path: str, target_lang: str, tone: str, openai_api_key: str, deepseek_api_key: str, use_deepseek=False, progress_callback=None, font_scale=1.0, use_smart_grouping=True, should_stop=None):
    """
    프레젠테이션을 번역하는 메인 함수
    progress_callback: (current_slide, total_slides, current_text) -> None
    """
    print(f"📂 파일: {pptx_path}")
    print(f"🌐 대상 언어: {target_lang}")
    print(f"🎙 톤: {tone}")
    if target_lang == "Chinese" and use_deepseek:
        print("🤖 DeepSeek 모델 사용 중...")
    else:
        print("🔑 OpenAI 클라이언트 초기화 중...")

    client = create_openai_client(openai_api_key)

    print("📖 프레젠테이션 로딩...")
    pres = Presentation(pptx_path)

    slide_count = len(pres.slides)
    print(f"🖼 슬라이드 수: {slide_count}")
    
    # 원본 한국어 텍스트 백업 (중국어 번역용)
    original_korean_backup = {}
    if "Chinese" in target_lang:
        chinese_type = "간체" if "Simplified" in target_lang else "번체"
        print(f"📝 원본 한국어 텍스트 백업 중... (중국어 {chinese_type} 번역용)")
        for s_idx, slide in enumerate(pres.slides, start=1):
            original_korean_backup[s_idx] = {}
            for shape_idx, shape in enumerate(slide.shapes):
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    tf = shape.text_frame
                    for p_idx, p in enumerate(tf.paragraphs):
                        tagged, _ = tag_paragraph(p)
                        if tagged:
                            original_korean_backup[s_idx][f"{shape_idx}_{p_idx}"] = tagged
                elif getattr(shape, "has_table", False) and shape.has_table:
                    for row_idx, row in enumerate(shape.table.rows):
                        for cell_idx, cell in enumerate(row.cells):
                            if not getattr(cell, "text_frame", None):
                                continue
                            tf = cell.text_frame
                            for p_idx, p in enumerate(tf.paragraphs):
                                tagged, _ = tag_paragraph(p)
                                if tagged:
                                    original_korean_backup[s_idx][f"table_{row_idx}_{cell_idx}_{p_idx}"] = tagged
    
    print("-" * 60)

    for s_idx, slide in enumerate(pres.slides, start=1):
        print(f"▶ 슬라이드 {s_idx}/{slide_count}")
        if progress_callback:
            progress_callback(s_idx, slide_count, f"슬라이드 {s_idx} 처리 중...")

        # 중지 신호 확인
        if should_stop and should_stop():
            print(f"⏹️ 슬라이드 {s_idx}에서 번역이 중지되었습니다.")
            return None
        
        # 스마트 그룹핑 사용 여부 확인
        if use_smart_grouping:
            # 스마트 그룹핑 먼저 시도
            smart_grouping_success = apply_smart_grouping_to_slide(
                slide, target_lang, tone, client, use_deepseek, font_scale
            )
            
            if not smart_grouping_success:
                # 스마트 그룹핑 실패 시 기존 방식으로 처리
                print(f"   🔄 기존 방식으로 슬라이드 {s_idx} 처리 중...")
                if not process_nested_shapes(
                    slide.shapes, target_lang, tone, client, use_deepseek, 
                    font_scale, progress_callback, s_idx, slide_count, "", should_stop
                ):
                    print(f"⏹️ 슬라이드 {s_idx}에서 번역이 중지되었습니다.")
                    return None
        else:
            # 스마트 그룹핑 비활성화 시 기존 방식으로 처리
            print(f"   🔄 기존 방식으로 슬라이드 {s_idx} 처리 중...")
            if not process_nested_shapes(
                slide.shapes, target_lang, tone, client, use_deepseek, 
                font_scale, progress_callback, s_idx, slide_count, "", should_stop
            ):
                print(f"⏹️ 슬라이드 {s_idx}에서 번역이 중지되었습니다.")
                return None

    folder = os.path.dirname(pptx_path)
    stem, ext = os.path.splitext(os.path.basename(pptx_path))
    
    # 중국어 번역의 경우 톤 대신 중국어 타입을 사용
    if "Chinese" in target_lang:
        chinese_type = "Simplified" if "Simplified" in target_lang else "Traditional"
        outfile_name = f"{stem}_Chinese_{chinese_type}_AI번역완료{ext}"
    else:
        safe_tone = re.sub(r'[^0-9A-Za-z가-힣_()+-]', '', tone.replace(' ', ''))
        outfile_name = f"{stem}_{target_lang}_{safe_tone}_AI번역완료{ext}"
    
    outfile_path = os.path.join(folder, outfile_name)
    outfile_path = unique_path(outfile_path)

    print("-" * 60)
    print("💾 저장 중...")
    pres.save(outfile_path)
    print(f"✅ 번역 완료! 저장된 파일: {outfile_path}")
    
    return outfile_path


# ========== [스마트 그룹핑 기능] ==========

def get_text_box_metadata(text_box):
    """텍스트박스의 메타데이터 추출"""
    metadata = {
        "text": "",
        "position": {"left": 0, "top": 0, "width": 0, "height": 0},
        "style": {},
        "has_text_frame": False
    }
    
    # 위치 정보
    if hasattr(text_box, 'left'):
        metadata["position"]["left"] = text_box.left
    if hasattr(text_box, 'top'):
        metadata["position"]["top"] = text_box.top
    if hasattr(text_box, 'width'):
        metadata["position"]["width"] = text_box.width
    if hasattr(text_box, 'height'):
        metadata["position"]["height"] = text_box.height
    
    # 텍스트 및 스타일 정보
    if hasattr(text_box, 'text_frame') and text_box.text_frame:
        metadata["has_text_frame"] = True
        metadata["text"] = text_box.text_frame.text
        
        # 첫 번째 문단의 첫 번째 run의 스타일
        if text_box.text_frame.paragraphs and text_box.text_frame.paragraphs[0].runs:
            first_run = text_box.text_frame.paragraphs[0].runs[0]
            # 안전한 색상 정보 추출
            font_color = None
            if first_run.font.color:
                try:
                    if hasattr(first_run.font.color, 'rgb') and first_run.font.color.rgb:
                        font_color = str(first_run.font.color.rgb)
                    elif hasattr(first_run.font.color, 'theme_color'):
                        font_color = f"theme_color_{first_run.font.color.theme_color}"
                except (AttributeError, TypeError):
                    font_color = None
            
            metadata["style"] = {
                "font_name": first_run.font.name,
                "font_size": first_run.font.size.pt if first_run.font.size else None,
                "font_color": font_color,
                "bold": first_run.font.bold,
                "italic": first_run.font.italic,
                "underline": first_run.font.underline
            }
    
    return metadata


def ai_analyze_text_grouping(text_metadata_list, slide_context=""):
    """AI가 텍스트박스들을 분석하여 그룹핑 결정 (GPT-5 사용)"""
    
    # 텍스트 목록 생성
    texts = [meta["text"] for meta in text_metadata_list if meta["text"].strip()]
    
    if not texts:
        return []
    
    prompt = f"""
다음 슬라이드의 텍스트들을 분석하여 그룹핑해주세요:

텍스트들: {texts}
슬라이드 컨텍스트: {slide_context}

그룹핑 규칙:
1. 문맥상 연결된 텍스트들은 같은 그룹 (예: "안녕하세요" + "반갑습니다")
2. 독립적인 정보 단위는 개별 그룹 (예: "제품명", "가격")
3. 스타일링을 위한 의도적 분할은 개별 그룹 (예: "100" + "개", "$" + "50")
4. 숫자+단위, 통화+금액, 라벨+값 등은 분리
5. 완전한 문장의 일부인 경우만 그룹핑

그룹핑 결과를 JSON 형태로 출력:
[
    {{"group": 1, "text_indices": [0, 1, 2]}},
    {{"group": 2, "text_indices": [3]}},
    {{"group": 3, "text_indices": [4, 5]}}
]

텍스트 인덱스는 0부터 시작합니다.
"""

    try:
        from openai import OpenAI
        client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        
        response = client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1000,
            temperature=0.1
        )
        
        content = response.choices[0].message.content.strip()
        
        # JSON 파싱
        import json
        import re
        
        # JSON 부분만 추출
        json_match = re.search(r'\[.*\]', content, re.DOTALL)
        if json_match:
            json_str = json_match.group(0)
            groups = json.loads(json_str)
            return groups
        else:
            print(f"⚠️ AI 응답에서 JSON을 찾을 수 없습니다: {content}")
            return []
            
    except Exception as e:
        print(f"❌ AI 그룹핑 분석 중 오류: {e}")
        return []


def smart_split_translation(translated_text, num_parts):
    """번역된 텍스트를 자연스럽게 분할 (GPT-5 사용)"""
    
    if num_parts <= 1:
        return [translated_text]
    
    prompt = f"""
다음 번역된 텍스트를 {num_parts}개의 부분으로 자연스럽게 분할해주세요:

번역된 텍스트: "{translated_text}"
분할 개수: {num_parts}

각 부분이 자연스러운 문장이나 구문이 되도록 분할해주세요.
의미 단위를 고려하여 분할하세요.

분할 결과를 JSON 배열로 출력:
["첫 번째 부분", "두 번째 부분", "세 번째 부분"]
"""

    try:
        from openai import OpenAI
        client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        
        response = client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=500,
            temperature=0.1
        )
        
        content = response.choices[0].message.content.strip()
        
        # JSON 파싱
        import json
        import re
        
        # JSON 배열 부분만 추출
        json_match = re.search(r'\[.*\]', content, re.DOTALL)
        if json_match:
            json_str = json_match.group(0)
            split_texts = json.loads(json_str)
            
            # 개수가 맞지 않으면 수동으로 분할
            if len(split_texts) != num_parts:
                print(f"⚠️ AI 분할 결과 개수 불일치: {len(split_texts)} != {num_parts}")
                return manual_split_text(translated_text, num_parts)
            
            return split_texts
        else:
            print(f"⚠️ AI 응답에서 JSON을 찾을 수 없습니다: {content}")
            return manual_split_text(translated_text, num_parts)
            
    except Exception as e:
        print(f"❌ AI 텍스트 분할 중 오류: {e}")
        return manual_split_text(translated_text, num_parts)


def manual_split_text(text, num_parts):
    """수동으로 텍스트 분할 (AI 실패 시 백업)"""
    if num_parts <= 1:
        return [text]
    
    # 공백으로 분할
    words = text.split()
    if len(words) <= num_parts:
        return words + [""] * (num_parts - len(words))
    
    # 균등하게 분할
    words_per_part = len(words) // num_parts
    result = []
    
    for i in range(num_parts):
        start = i * words_per_part
        if i == num_parts - 1:  # 마지막 부분은 나머지 모든 단어
            end = len(words)
        else:
            end = (i + 1) * words_per_part
        
        part = " ".join(words[start:end])
        result.append(part)
    
    return result


def translate_text_group_with_style_preservation(group_metadata, target_lang, tone, client, use_deepseek=False):
    """텍스트 그룹을 번역하되 각 텍스트박스의 스타일 보존"""
    
    if not group_metadata:
        return []
    
    # 1. 그룹의 모든 텍스트를 하나로 합침
    combined_text = " ".join([meta["text"] for meta in group_metadata if meta["text"].strip()])
    
    if not combined_text.strip():
        return group_metadata
    
    # 2. 전체 문맥으로 번역
    print(f"   🔤 그룹 번역 중: {combined_text[:50]}...")
    
    # 태그된 텍스트로 변환 (기존 로직 활용)
    tagged_text = f"[[R1]]{combined_text}[[/R1]]"
    translated = gpt_translate_tagged(tagged_text, client, target_lang, tone, use_deepseek)
    translated = translated.strip().strip('"').strip("'")
    
    # 3. 번역된 텍스트를 원래 텍스트박스 개수만큼 분할
    num_parts = len(group_metadata)
    split_texts = smart_split_translation(translated, num_parts)
    
    # 4. 각 텍스트박스에 번역된 텍스트와 원래 스타일 적용
    result = []
    for i, box_meta in enumerate(group_metadata):
        translated_text = split_texts[i] if i < len(split_texts) else ""
        
        result.append({
            "text": translated_text,
            "position": box_meta["position"],
            "style": box_meta["style"],
            "has_text_frame": box_meta["has_text_frame"]
        })
    
    return result


def apply_smart_grouping_to_slide(slide, target_lang, tone, client, use_deepseek=False, font_scale=1.0):
    """슬라이드에 스마트 그룹핑 적용"""
    
    print(f"   🧠 스마트 그룹핑 적용 중...")
    
    # 1. 모든 텍스트박스 수집 및 메타데이터 추출
    text_boxes = []
    text_metadata = []
    
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text.strip():
            text_boxes.append(shape)
            metadata = get_text_box_metadata(shape)
            text_metadata.append(metadata)
    
    if not text_metadata:
        return True
    
    # 2. AI 기반 그룹핑 분석
    groups = ai_analyze_text_grouping(text_metadata, f"슬라이드 {slide.slide_id}")
    
    if not groups:
        print(f"   ⚠️ AI 그룹핑 실패, 개별 번역으로 진행")
        return False
    
    # 3. 각 그룹 번역 및 적용
    for group_info in groups:
        group_indices = group_info.get("text_indices", [])
        if not group_indices:
            continue
        
        # 그룹의 메타데이터 수집
        group_metadata = [text_metadata[i] for i in group_indices if i < len(text_metadata)]
        
        if len(group_metadata) <= 1:
            # 단일 텍스트박스는 기존 방식으로 처리
            continue
        
        # 그룹 번역
        translated_group = translate_text_group_with_style_preservation(
            group_metadata, target_lang, tone, client, use_deepseek
        )
        
        # 원래 텍스트박스에 적용
        for i, box_idx in enumerate(group_indices):
            if i < len(translated_group) and box_idx < len(text_boxes):
                original_box = text_boxes[box_idx]
                translated_text = translated_group[i]["text"]
                
                # 텍스트 업데이트
                update_text_box_with_translation(original_box, translated_text, font_scale)
    
    print(f"   ✅ 스마트 그룹핑 완료")
    return True


def update_text_box_with_translation(text_box, translated_text, font_scale=1.0):
    """텍스트박스에 번역된 텍스트 적용"""
    
    if not hasattr(text_box, 'text_frame') or not text_box.text_frame:
        return
    
    # 기존 내용을 번역된 텍스트로 교체
    for paragraph in text_box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = ""
    
    # 첫 번째 문단의 첫 번째 run에 번역된 텍스트 적용
    if text_box.text_frame.paragraphs and text_box.text_frame.paragraphs[0].runs:
        first_run = text_box.text_frame.paragraphs[0].runs[0]
        first_run.text = translated_text
        
        # 폰트 크기 조정
        if font_scale != 1.0 and first_run.font.size:
            from pptx.util import Pt
            original_size = first_run.font.size.pt
            adjusted_size = original_size * font_scale
            first_run.font.size = Pt(adjusted_size)
