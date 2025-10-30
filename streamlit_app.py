# streamlit_app.py
# Main Streamlit application for PPT translation and comparison

import streamlit as st
import tempfile
import os
import time
from typing import Dict, List, Optional
import traceback

# Import our custom modules
from translation_engine import (
    translate_presentation, LANG_OPTIONS, TONE_OPTIONS, 
    create_openai_client, create_deepseek_client
)
from accuracy_checker import (
    evaluate_translation_quality, batch_evaluate_translations,
    get_flagged_translations, get_low_confidence_translations
)
from comparison_ui import (
    create_slide_comparison_ui, create_flagged_translations_ui,
    create_retranslation_ui, extract_slide_text_data
)


# Page configuration
st.set_page_config(
    page_title="PPT AI Translation & Comparison",
    page_icon="🌐",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        font-weight: bold;
        color: #1f77b4;
        text-align: center;
        margin-bottom: 2rem;
    }
    .sub-header {
        font-size: 1.5rem;
        font-weight: bold;
        color: #2c3e50;
        margin-top: 1.5rem;
        margin-bottom: 1rem;
    }
    .success-box {
        background-color: #d4edda;
        border: 1px solid #c3e6cb;
        border-radius: 0.375rem;
        padding: 1rem;
        margin: 1rem 0;
    }
    .warning-box {
        background-color: #fff3cd;
        border: 1px solid #ffeaa7;
        border-radius: 0.375rem;
        padding: 1rem;
        margin: 1rem 0;
    }
    .error-box {
        background-color: #f8d7da;
        border: 1px solid #f5c6cb;
        border-radius: 0.375rem;
        padding: 1rem;
        margin: 1rem 0;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state
if "current_slide" not in st.session_state:
    st.session_state.current_slide = 1
if "translation_results" not in st.session_state:
    st.session_state.translation_results = None
if "evaluation_results" not in st.session_state:
    st.session_state.evaluation_results = []
if "flagged_translations" not in st.session_state:
    st.session_state.flagged_translations = []
if "retranslate_data" not in st.session_state:
    st.session_state.retranslate_data = None
if "is_translating" not in st.session_state:
    st.session_state.is_translating = False
if "translation_cancelled" not in st.session_state:
    st.session_state.translation_cancelled = False


def main():
    """Main application function"""
    
    # Header
    st.markdown('<h1 class="main-header">🌐 PPT AI Translation & Comparison</h1>', unsafe_allow_html=True)
    
    # Get API keys from secrets (hidden from users)
    openai_api_key = st.secrets.get("OPENAI_API_KEY", "")
    deepseek_api_key = st.secrets.get("DEEPSEEK_API_KEY", "")
    
    # Sidebar for settings
    with st.sidebar:
        st.header("⚙️ 번역 설정")
        
        # Translation Settings
        target_language = st.selectbox(
            "대상 언어",
            LANG_OPTIONS,
            index=0,
            help="번역할 언어를 선택하세요"
        )
        
        tone = st.selectbox(
            "톤 선택",
            TONE_OPTIONS,
            index=0,
            help="번역 톤을 선택하세요"
        )
        
        use_deepseek = st.checkbox(
            "중국어 번역 시 DeepSeek 사용",
            value=True,
            disabled="Chinese" not in target_language,
            help="중국어 번역 시 DeepSeek API를 사용합니다"
        )
        
        # Font size adjustment
        st.markdown("---")
        st.subheader("🎨 폰트 설정")
        
        font_scale = st.slider(
            "폰트 크기 배율",
            min_value=0.5,
            max_value=2.0,
            value=1.0,
            step=0.1,
            help="1.0이 기본 크기입니다. 0.5는 절반 크기, 2.0은 두 배 크기입니다"
        )
        
        # Show preview
        if font_scale != 1.0:
            if font_scale < 1.0:
                st.info(f"폰트 크기가 {font_scale:.1f}배로 조정됩니다 (작아짐)")
            else:
                st.info(f"폰트 크기가 {font_scale:.1f}배로 조정됩니다 (커짐)")
        
        # Fixed accuracy settings (hidden from users)
        confidence_threshold = 70
        auto_evaluate = True
    
    # Main content area
    tab1, tab2, tab3, tab4 = st.tabs(["📤 번역", "📊 비교", "🚩 검토", "ℹ️ 정보"])
    
    with tab1:
        translation_tab(openai_api_key, deepseek_api_key, target_language, tone, use_deepseek, auto_evaluate, confidence_threshold, font_scale)
    
    with tab2:
        comparison_tab()
    
    with tab3:
        review_tab(openai_api_key, deepseek_api_key, target_language, tone, use_deepseek)
    
    with tab4:
        info_tab()


def translation_tab(openai_api_key: str, deepseek_api_key: str, target_language: str, 
                   tone: str, use_deepseek: bool, auto_evaluate: bool, confidence_threshold: int, font_scale: float = 1.0):
    """Translation tab content"""
    
    st.markdown('<h2 class="sub-header">📤 PPT 번역</h2>', unsafe_allow_html=True)
    
    # File upload
    uploaded_file = st.file_uploader(
        "PPTX 파일을 업로드하세요",
        type=['pptx'],
        help="번역할 PowerPoint 파일을 선택하세요"
    )
    
    if uploaded_file is not None:
        # Display file info
        st.success(f"✅ 파일 업로드 완료: {uploaded_file.name}")
        st.write(f"파일 크기: {uploaded_file.size:,} bytes")
        
        # API key validation
        if not openai_api_key:
            st.error("❌ 서비스가 일시적으로 사용할 수 없습니다. 잠시 후 다시 시도해주세요.")
            return
        
        if "Chinese" in target_language and use_deepseek and not deepseek_api_key:
            st.error("❌ 중국어 번역 서비스가 일시적으로 사용할 수 없습니다. 다른 언어를 선택하거나 잠시 후 다시 시도해주세요.")
            return
        
        # Translation button with start/stop functionality
        col1, col2 = st.columns([1, 1])
        
        with col1:
            if not st.session_state.is_translating:
                if st.button("🚀 번역 시작", type="primary", use_container_width=True):
                    st.session_state.is_translating = True
                    st.session_state.translation_cancelled = False
                    st.rerun()
            else:
                if st.button("⏹️ 번역 중지", type="secondary", use_container_width=True):
                    st.session_state.translation_cancelled = True
                    st.session_state.is_translating = False
                    st.rerun()
        
        with col2:
            if st.session_state.is_translating:
                st.info("🔄 번역 진행 중...")
        
        # Start translation if button was clicked
        if st.session_state.is_translating and not st.session_state.translation_cancelled:
            # Create a stop signal function
            def should_stop():
                return st.session_state.translation_cancelled
            
            translate_file(uploaded_file, openai_api_key, deepseek_api_key, 
                          target_language, tone, use_deepseek, auto_evaluate, confidence_threshold, font_scale, should_stop)


def translate_file(uploaded_file, openai_api_key: str, deepseek_api_key: str, 
                  target_language: str, tone: str, use_deepseek: bool, 
                  auto_evaluate: bool, confidence_threshold: int, font_scale: float = 1.0, should_stop=None):
    """Handle file translation"""
    
    try:
        # Save uploaded file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_file_path = tmp_file.name
        
        # Progress tracking
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        def progress_callback(current_slide, total_slides, current_text):
            if total_slides > 0:
                progress = min(current_slide / total_slides, 1.0)  # 1.0을 초과하지 않도록 제한
            else:
                progress = 0.0
            progress_bar.progress(progress)
            if total_slides > 0:
                status_text.text(f"진행률: {current_slide}/{total_slides} - {current_text}")
            else:
                status_text.text(f"처리 중: {current_text}")
        
        # Start translation
        with st.spinner("번역 중..."):
            translated_file_path = translate_presentation(
                tmp_file_path,
                target_language,
                tone,
                openai_api_key,
                deepseek_api_key,
                use_deepseek,
                progress_callback,
                font_scale,
                should_stop
            )
        
        # Check if translation was cancelled
        if translated_file_path is None:
            st.warning("⏹️ 번역이 중지되었습니다.")
            st.session_state.is_translating = False
            st.session_state.translation_cancelled = False
            return
        
        # Translation completed
        progress_bar.progress(1.0)
        status_text.text("✅ 번역 완료!")
        
        # Store results in session state
        st.session_state.translation_results = {
            "original_file": tmp_file_path,
            "translated_file": translated_file_path,
            "target_language": target_language,
            "tone": tone,
            "use_deepseek": use_deepseek
        }
        
        # Auto-evaluation if enabled
        if auto_evaluate:
            with st.spinner("정확도 평가 중..."):
                evaluate_translations(translated_file_path, openai_api_key, deepseek_api_key, 
                                    target_language, use_deepseek, confidence_threshold)
        
        # Download button
        with open(translated_file_path, "rb") as file:
            st.download_button(
                label="📥 번역된 파일 다운로드",
                data=file.read(),
                file_name=os.path.basename(translated_file_path),
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )
        
        st.success("🎉 번역이 성공적으로 완료되었습니다!")
        
        # Reset translation state
        st.session_state.is_translating = False
        st.session_state.translation_cancelled = False
        
        # Show next steps
        st.info("💡 다음 단계: '비교' 탭에서 원본과 번역본을 비교하거나, '검토' 탭에서 플래그된 번역을 확인하세요.")
        
    except Exception as e:
        st.error(f"❌ 번역 중 오류가 발생했습니다: {str(e)}")
        st.error(f"상세 오류: {traceback.format_exc()}")
        
        # Reset translation state on error
        st.session_state.is_translating = False
        st.session_state.translation_cancelled = False
    
    finally:
        # Clean up temporary file
        try:
            if 'tmp_file_path' in locals():
                os.unlink(tmp_file_path)
        except:
            pass


def evaluate_translations(translated_file_path: str, openai_api_key: str, deepseek_api_key: str,
                         target_language: str, use_deepseek: bool, confidence_threshold: int):
    """Evaluate translation quality"""
    
    try:
        # Extract translation data for evaluation
        translations_data = extract_translation_data_for_evaluation(translated_file_path)
        
        if not translations_data:
            st.warning("평가할 번역 데이터가 없습니다.")
            return
        
        # Batch evaluate translations
        evaluation_results = batch_evaluate_translations(
            translations_data,
            target_language,
            openai_api_key,
            deepseek_api_key,
            use_deepseek
        )
        
        # Store results
        st.session_state.evaluation_results = evaluation_results
        st.session_state.flagged_translations = get_flagged_translations(evaluation_results)
        
        # Show summary
        total_translations = len(evaluation_results)
        flagged_count = len(st.session_state.flagged_translations)
        low_confidence_count = len(get_low_confidence_translations(evaluation_results, confidence_threshold))
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.metric("총 번역 수", total_translations)
        
        with col2:
            st.metric("플래그된 번역", flagged_count, delta=f"{flagged_count/total_translations*100:.1f}%" if total_translations > 0 else "0%")
        
        with col3:
            st.metric("낮은 신뢰도", low_confidence_count, delta=f"{low_confidence_count/total_translations*100:.1f}%" if total_translations > 0 else "0%")
        
        if flagged_count > 0:
            st.warning(f"⚠️ {flagged_count}개의 번역이 검토가 필요합니다. '검토' 탭을 확인하세요.")
        else:
            st.success("✅ 모든 번역이 양호한 품질입니다!")
            
    except Exception as e:
        st.error(f"정확도 평가 중 오류가 발생했습니다: {str(e)}")


def extract_translation_data_for_evaluation(pptx_path: str) -> List[Dict]:
    """Extract translation data for evaluation"""
    
    try:
        from pptx import Presentation
        pres = Presentation(pptx_path)
        translations_data = []
        
        for slide_idx, slide in enumerate(pres.slides, 1):
            for shape_idx, shape in enumerate(slide.shapes):
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    tf = shape.text_frame
                    for p_idx, paragraph in enumerate(tf.paragraphs):
                        if paragraph.text.strip():
                            translations_data.append({
                                "id": f"slide_{slide_idx}_shape_{shape_idx}_para_{p_idx}",
                                "slide_number": slide_idx,
                                "shape_type": "text_frame",
                                "shape_idx": shape_idx,
                                "paragraph_idx": p_idx,
                                "original_text": paragraph.text.strip(),  # This would need to be stored separately
                                "translated_text": paragraph.text.strip()
                            })
                elif getattr(shape, "has_table", False) and shape.has_table:
                    for row_idx, row in enumerate(shape.table.rows):
                        for cell_idx, cell in enumerate(row.cells):
                            if getattr(cell, "text_frame", None):
                                tf = cell.text_frame
                                for p_idx, paragraph in enumerate(tf.paragraphs):
                                    if paragraph.text.strip():
                                        translations_data.append({
                                            "id": f"slide_{slide_idx}_table_{shape_idx}_row_{row_idx}_cell_{cell_idx}_para_{p_idx}",
                                            "slide_number": slide_idx,
                                            "shape_type": "table_cell",
                                            "shape_idx": shape_idx,
                                            "row_idx": row_idx,
                                            "cell_idx": cell_idx,
                                            "paragraph_idx": p_idx,
                                            "original_text": paragraph.text.strip(),  # This would need to be stored separately
                                            "translated_text": paragraph.text.strip()
                                        })
        
        return translations_data
        
    except Exception as e:
        st.error(f"번역 데이터 추출 중 오류: {str(e)}")
        return []


def comparison_tab():
    """Comparison tab content"""
    
    st.markdown('<h2 class="sub-header">📊 원본 vs 번역본 비교</h2>', unsafe_allow_html=True)
    
    if not st.session_state.translation_results:
        st.info("먼저 번역을 완료해주세요.")
        return
    
    # Get file paths
    original_file = st.session_state.translation_results["original_file"]
    translated_file = st.session_state.translation_results["translated_file"]
    
    if not os.path.exists(original_file) or not os.path.exists(translated_file):
        st.error("번역된 파일을 찾을 수 없습니다. 다시 번역해주세요.")
        return
    
    # Create comparison UI
    create_slide_comparison_ui(
        original_file,
        translated_file,
        st.session_state.evaluation_results,
        st.session_state.current_slide
    )


def review_tab(openai_api_key: str, deepseek_api_key: str, target_language: str, tone: str, use_deepseek: bool):
    """Review tab content"""
    
    st.markdown('<h2 class="sub-header">🚩 플래그된 번역 검토</h2>', unsafe_allow_html=True)
    
    if not st.session_state.flagged_translations:
        if st.session_state.evaluation_results:
            st.success("✅ 검토가 필요한 번역이 없습니다!")
        else:
            st.info("먼저 번역을 완료하고 정확도 평가를 실행해주세요.")
        return
    
    # Show flagged translations
    create_flagged_translations_ui(st.session_state.flagged_translations)
    
    # Retranslation UI if needed
    if st.session_state.retranslate_data:
        st.markdown("---")
        create_retranslation_ui(
            st.session_state.retranslate_data,
            openai_api_key,
            deepseek_api_key,
            target_language,
            tone,
            use_deepseek
        )


def info_tab():
    """Info tab content"""
    
    st.markdown('<h2 class="sub-header">ℹ️ 정보</h2>', unsafe_allow_html=True)
    
    st.markdown("""
    ## 🌐 PPT AI Translation & Comparison
    
    이 애플리케이션은 PowerPoint 프레젠테이션을 AI를 사용하여 번역하고, 번역 품질을 평가하며, 원본과 번역본을 비교할 수 있는 도구입니다.
    
    ### 주요 기능
    
    #### 📤 번역
    - **다국어 지원**: 12개 언어로 번역 가능
    - **톤 선택**: 다양한 비즈니스 톤으로 번역
    - **서식 보존**: 원본 PPT의 서식과 레이아웃 유지
    - **API 선택**: OpenAI GPT-4 또는 DeepSeek 사용 가능
    
    #### 📊 비교
    - **슬라이드별 비교**: 원본과 번역본을 슬라이드별로 비교
    - **텍스트 비교**: 상세한 텍스트 내용 비교
    - **시각적 비교**: 슬라이드 이미지 비교 (개발 중)
    
    #### 🚩 검토
    - **자동 품질 평가**: AI가 번역 품질을 자동으로 평가
    - **신뢰도 점수**: 0-100% 신뢰도 점수 제공
    - **문제점 감지**: 번역의 문제점을 자동으로 감지
    - **재번역**: 문제가 있는 번역을 쉽게 재번역
    
    ### 사용 방법
    
    1. **API 키 설정**: 사이드바에서 OpenAI API 키를 입력하세요
    2. **번역 설정**: 대상 언어와 톤을 선택하세요
    3. **파일 업로드**: 번역할 PPTX 파일을 업로드하세요
    4. **번역 실행**: "번역 시작" 버튼을 클릭하세요
    5. **결과 확인**: 비교 탭에서 원본과 번역본을 비교하세요
    6. **품질 검토**: 검토 탭에서 플래그된 번역을 확인하세요
    
    ### 지원 언어
    
    - English
    - Indonesian
    - Italian
    - French
    - Spanish
    - Korean
    - Japanese
    - Russian
    - German
    - Portuguese
    - Chinese (Simplified)
    - Chinese (Traditional)
    
    ### 톤 옵션
    
    - **기본값**: 일반적인 뷰티 업계 톤
    - **Med/Pharma Pro (20y)**: 의료기기/전문약사 20년 전문가 톤
    - **Medical/Science Expert**: 의학/과학 번역 전문가 (정확한 의학 용어 사용)
    - **Beauty Pro (20y, chic)**: 세련된 뷰티 20년 전문가 톤
    - **GenZ Female (20s)**: 20대 여성 타깃의 친근한 톤
    
    ### 기술 스택
    
    - **Frontend**: Streamlit
    - **AI Models**: OpenAI GPT-4, DeepSeek
    - **PPT Processing**: python-pptx
    - **Deployment**: Streamlit Cloud
    
    ### 문제 해결
    
    **Q: 번역이 제대로 되지 않아요**
    A: API 키가 올바른지 확인하고, 네트워크 연결을 확인해주세요.
    
    **Q: 서식이 깨져요**
    A: 복잡한 서식이 있는 경우 일부 서식이 변경될 수 있습니다. 텍스트 내용은 정확히 번역됩니다.
    
    **Q: 신뢰도 점수가 낮아요**
    A: 낮은 신뢰도는 번역 품질이 개선될 수 있음을 의미합니다. 재번역 기능을 사용해보세요.
    
    ### 지원
    
    문제가 있거나 개선 사항이 있으시면 GitHub 저장소에 이슈를 등록해주세요.
    """)
    
    # GitHub link
    st.markdown("### 🔗 링크")
    st.markdown("[GitHub 저장소](https://github.com/Teletobimy/ppt_ai_translation)")


if __name__ == "__main__":
    main()
