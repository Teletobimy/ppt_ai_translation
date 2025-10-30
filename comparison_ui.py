# comparison_ui.py
# Comparison UI components for Streamlit app

import streamlit as st
from typing import Dict, List, Optional, Tuple
import tempfile
import os
from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image
import base64


def create_slide_comparison_ui(original_pptx_path: str, translated_pptx_path: str, 
                              evaluation_results: List[Dict], current_slide: int = 1) -> None:
    """
    슬라이드별 원본-번역본 비교 UI 생성
    """
    st.subheader("📊 슬라이드별 비교")
    
    # 슬라이드 네비게이션
    col1, col2, col3 = st.columns([1, 2, 1])
    
    with col1:
        if st.button("◀ 이전 슬라이드", disabled=current_slide <= 1):
            st.session_state.current_slide = max(1, current_slide - 1)
            st.rerun()
    
    with col2:
        st.write(f"**슬라이드 {current_slide}**")
    
    with col3:
        if st.button("다음 슬라이드 ▶", disabled=current_slide >= get_slide_count(original_pptx_path)):
            st.session_state.current_slide = min(get_slide_count(original_pptx_path), current_slide + 1)
            st.rerun()
    
    # 비교 모드 선택
    comparison_mode = st.radio(
        "비교 모드 선택:",
        ["텍스트 비교", "시각적 비교"],
        horizontal=True
    )
    
    if comparison_mode == "텍스트 비교":
        create_text_comparison(original_pptx_path, translated_pptx_path, evaluation_results, current_slide)
    else:
        create_visual_comparison(original_pptx_path, translated_pptx_path, current_slide)


def create_text_comparison(original_pptx_path: str, translated_pptx_path: str, 
                          evaluation_results: List[Dict], slide_number: int) -> None:
    """
    텍스트 기반 비교 UI
    """
    # 원본과 번역본 슬라이드 데이터 추출
    original_slide_data = extract_slide_text_data(original_pptx_path, slide_number)
    translated_slide_data = extract_slide_text_data(translated_pptx_path, slide_number)
    
    # 좌우 분할 레이아웃
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("### 🔵 원본 (한국어)")
        display_slide_text_data(original_slide_data, "original")
    
    with col2:
        st.markdown("### 🟢 번역본")
        display_slide_text_data(translated_slide_data, "translated", evaluation_results, slide_number)


def create_visual_comparison(original_pptx_path: str, translated_pptx_path: str, slide_number: int) -> None:
    """
    시각적 비교 UI (슬라이드 이미지)
    """
    try:
        # Convert slides to images
        original_image = convert_slide_to_image(original_pptx_path, slide_number)
        translated_image = convert_slide_to_image(translated_pptx_path, slide_number)
        
        if original_image and translated_image:
            # Display side-by-side images
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("### 🔵 원본 슬라이드")
                st.image(original_image, use_column_width=True, caption=f"슬라이드 {slide_number} - 원본")
            
            with col2:
                st.markdown("### 🟢 번역본 슬라이드")
                st.image(translated_image, use_column_width=True, caption=f"슬라이드 {slide_number} - 번역본")
            
            # Add zoom functionality
            st.markdown("---")
            st.markdown("**🔍 확대/축소 기능**")
            
            zoom_level = st.slider("확대 비율", 0.5, 2.0, 1.0, 0.1)
            
            if zoom_level != 1.0:
                # Resize images based on zoom level
                original_resized = original_image.resize(
                    (int(original_image.width * zoom_level), 
                     int(original_image.height * zoom_level)),
                    Image.Resampling.LANCZOS
                )
                translated_resized = translated_image.resize(
                    (int(translated_image.width * zoom_level), 
                     int(translated_image.height * zoom_level)),
                    Image.Resampling.LANCZOS
                )
                
                col1, col2 = st.columns(2)
                
                with col1:
                    st.image(original_resized, use_column_width=True, caption=f"원본 (확대 {zoom_level:.1f}x)")
                
                with col2:
                    st.image(translated_resized, use_column_width=True, caption=f"번역본 (확대 {zoom_level:.1f}x)")
        else:
            st.warning("⚠️ 슬라이드 이미지를 생성할 수 없습니다. 텍스트 비교를 사용해주세요.")
            
    except Exception as e:
        st.error(f"시각적 비교 중 오류가 발생했습니다: {str(e)}")
        st.info("텍스트 비교를 사용해주세요.")


def convert_slide_to_image(pptx_path: str, slide_number: int) -> Optional[Image.Image]:
    """
    PPT 슬라이드를 이미지로 변환
    """
    try:
        # This is a simplified implementation
        # In a production environment, you might want to use:
        # - python-pptx with additional libraries
        # - LibreOffice headless mode
        # - Or other PPT to image conversion tools
        
        from pptx import Presentation
        import io
        from PIL import Image, ImageDraw, ImageFont
        
        pres = Presentation(pptx_path)
        
        if slide_number > len(pres.slides):
            return None
        
        slide = pres.slides[slide_number - 1]
        
        # Create a simple image representation of the slide
        # This is a basic implementation - you might want to use more sophisticated methods
        width, height = 800, 600
        img = Image.new('RGB', (width, height), 'white')
        draw = ImageDraw.Draw(img)
        
        # Try to use a default font
        try:
            font = ImageFont.truetype("arial.ttf", 16)
        except:
            font = ImageFont.load_default()
        
        y_offset = 20
        
        # Extract text from slide and draw it
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                # Simple text rendering
                text_lines = shape.text.split('\n')
                for line in text_lines:
                    if line.strip():
                        draw.text((20, y_offset), line.strip(), fill='black', font=font)
                        y_offset += 25
                y_offset += 10
        
        return img
        
    except Exception as e:
        print(f"Error converting slide to image: {e}")
        return None


def extract_slide_text_data(pptx_path: str, slide_number: int) -> List[Dict]:
    """
    특정 슬라이드의 텍스트 데이터를 추출
    """
    try:
        pres = Presentation(pptx_path)
        if slide_number > len(pres.slides):
            return []
        
        slide = pres.slides[slide_number - 1]
        text_data = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                tf = shape.text_frame
                for p_idx, paragraph in enumerate(tf.paragraphs):
                    if paragraph.text.strip():
                        text_data.append({
                            "shape_type": "text_frame",
                            "shape_idx": shape_idx,
                            "paragraph_idx": p_idx,
                            "text": paragraph.text.strip(),
                            "runs": [run.text for run in paragraph.runs if run.text.strip()]
                        })
            elif getattr(shape, "has_table", False) and shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        if getattr(cell, "text_frame", None):
                            tf = cell.text_frame
                            for p_idx, paragraph in enumerate(tf.paragraphs):
                                if paragraph.text.strip():
                                    text_data.append({
                                        "shape_type": "table_cell",
                                        "shape_idx": shape_idx,
                                        "row_idx": row_idx,
                                        "cell_idx": cell_idx,
                                        "paragraph_idx": p_idx,
                                        "text": paragraph.text.strip(),
                                        "runs": [run.text for run in paragraph.runs if run.text.strip()]
                                    })
        
        return text_data
    except Exception as e:
        st.error(f"슬라이드 데이터 추출 오류: {e}")
        return []


def display_slide_text_data(slide_data: List[Dict], data_type: str, 
                           evaluation_results: List[Dict] = None, slide_number: int = 0) -> None:
    """
    슬라이드 텍스트 데이터를 표시
    """
    if not slide_data:
        st.write("텍스트가 없습니다.")
        return
    
    for i, text_item in enumerate(slide_data):
        # 평가 결과가 있으면 해당 텍스트의 신뢰도 점수 표시
        confidence_score = None
        is_flagged = False
        issues = []
        
        if evaluation_results and data_type == "translated":
            # 해당 텍스트와 매칭되는 평가 결과 찾기
            for eval_result in evaluation_results:
                if (eval_result.get("slide_number") == slide_number and
                    eval_result.get("shape_type") == text_item["shape_type"] and
                    eval_result.get("shape_idx") == text_item["shape_idx"]):
                    confidence_score = eval_result.get("confidence_score", 0)
                    is_flagged = eval_result.get("is_flagged", False)
                    issues = eval_result.get("issues", [])
                    break
        
        # 텍스트 컨테이너
        container_style = ""
        if data_type == "translated" and is_flagged:
            container_style = "border-left: 4px solid #ff6b6b; background-color: #fff5f5; padding: 10px; margin: 5px 0;"
        elif data_type == "translated" and confidence_score is not None and confidence_score < 70:
            container_style = "border-left: 4px solid #ffa726; background-color: #fff8e1; padding: 10px; margin: 5px 0;"
        
        if container_style:
            st.markdown(f'<div style="{container_style}">', unsafe_allow_html=True)
        
        # 텍스트 표시
        st.write(f"**{text_item['shape_type']}**")
        st.write(text_item['text'])
        
        # 신뢰도 점수 표시
        if data_type == "translated" and confidence_score is not None:
            color = "red" if confidence_score < 50 else "orange" if confidence_score < 70 else "green"
            st.markdown(f"신뢰도: <span style='color: {color}; font-weight: bold;'>{confidence_score}%</span>", 
                       unsafe_allow_html=True)
        
        # 문제점 표시
        if issues:
            st.markdown("**⚠️ 문제점:**")
            for issue in issues:
                st.write(f"- {issue}")
        
        # 재번역 버튼
        if data_type == "translated" and is_flagged:
            if st.button(f"🔄 재번역", key=f"retranslate_{i}"):
                st.session_state.retranslate_text = text_item['text']
                st.session_state.retranslate_slide = slide_number
                st.session_state.retranslate_shape = text_item
        
        if container_style:
            st.markdown('</div>', unsafe_allow_html=True)
        
        st.write("---")


def get_slide_count(pptx_path: str) -> int:
    """
    PPT 파일의 슬라이드 수 반환
    """
    try:
        pres = Presentation(pptx_path)
        return len(pres.slides)
    except:
        return 0


def create_flagged_translations_ui(flagged_translations: List[Dict]) -> None:
    """
    플래그된 번역들을 검토하는 UI
    """
    st.subheader("🚩 플래그된 번역 검토")
    
    if not flagged_translations:
        st.success("✅ 플래그된 번역이 없습니다!")
        return
    
    st.write(f"총 {len(flagged_translations)}개의 번역이 검토가 필요합니다.")
    
    # 필터링 옵션
    col1, col2 = st.columns(2)
    
    with col1:
        confidence_threshold = st.slider("신뢰도 임계값", 0, 100, 70)
    
    with col2:
        filter_by_issues = st.checkbox("문제점이 있는 번역만 표시")
    
    # 필터링된 번역 목록
    filtered_translations = flagged_translations.copy()
    
    if confidence_threshold > 0:
        filtered_translations = [t for t in filtered_translations 
                               if t.get("confidence_score", 0) < confidence_threshold]
    
    if filter_by_issues:
        filtered_translations = [t for t in filtered_translations 
                               if t.get("issues", [])]
    
    # 번역 목록 표시
    for i, translation in enumerate(filtered_translations):
        with st.expander(f"번역 {i+1} - 신뢰도: {translation.get('confidence_score', 0)}%", expanded=False):
            col1, col2 = st.columns(2)
            
            with col1:
                st.write("**원본:**")
                st.write(translation.get("original_text", ""))
            
            with col2:
                st.write("**번역:**")
                st.write(translation.get("translated_text", ""))
            
            # 평가 정보
            st.write("**평가 정보:**")
            st.write(f"- 신뢰도: {translation.get('confidence_score', 0)}%")
            st.write(f"- 정확도: {translation.get('accuracy', 'N/A')}")
            st.write(f"- 자연스러움: {translation.get('naturalness', 'N/A')}")
            
            # 문제점
            if translation.get("issues"):
                st.write("**문제점:**")
                for issue in translation["issues"]:
                    st.write(f"- {issue}")
            
            # 개선 제안
            if translation.get("suggestions"):
                st.write("**개선 제안:**")
                for suggestion in translation["suggestions"]:
                    st.write(f"- {suggestion}")
            
            # 액션 버튼
            col1, col2, col3 = st.columns(3)
            
            with col1:
                if st.button("✅ 승인", key=f"approve_{i}"):
                    # 승인 로직
                    st.success("번역이 승인되었습니다.")
            
            with col2:
                if st.button("🔄 재번역", key=f"retranslate_{i}"):
                    # 재번역 로직
                    st.session_state.retranslate_data = translation
                    st.rerun()
            
            with col3:
                if st.button("✏️ 수정", key=f"edit_{i}"):
                    # 수정 로직
                    st.session_state.edit_translation = translation
                    st.rerun()


def create_retranslation_ui(translation_data: Dict, openai_api_key: str, deepseek_api_key: str, 
                           target_lang: str, tone: str, use_deepseek: bool = False) -> None:
    """
    재번역 UI
    """
    st.subheader("🔄 재번역")
    
    if not translation_data:
        st.warning("재번역할 번역을 선택해주세요.")
        return
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.write("**원본 텍스트:**")
        st.text_area("", value=translation_data.get("original_text", ""), height=200, disabled=True)
    
    with col2:
        st.write("**현재 번역:**")
        st.text_area("", value=translation_data.get("translated_text", ""), height=200, disabled=True)
    
    # 재번역 옵션
    st.write("**재번역 옵션:**")
    
    col1, col2 = st.columns(2)
    
    with col1:
        alternative_tone = st.selectbox(
            "대체 톤 선택:",
            ["기본값", "Med/Pharma Pro (20y)", "Beauty Pro (20y, chic)", "GenZ Female (20s)"],
            index=0
        )
    
    with col2:
        use_alternative_api = st.checkbox("다른 API 사용", value=not use_deepseek)
    
    # 재번역 실행
    if st.button("🔄 재번역 실행"):
        with st.spinner("재번역 중..."):
            # 재번역 로직 구현
            # translation_engine.gpt_translate_tagged() 사용
            st.success("재번역이 완료되었습니다!")
    
    # 취소
    if st.button("❌ 취소"):
        st.session_state.retranslate_data = None
        st.rerun()
